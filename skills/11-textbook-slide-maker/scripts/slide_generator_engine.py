import os
import sys
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

sys.stdout.reconfigure(encoding='utf-8')

# ==============================================================================
# 金牌学术课件设计规范 (Benchmark: 《第十四章_统筹发展和安全.pptx》)
# textbook-slide-maker 核心引擎 v1.6 (讲练融合 + 图表嵌入 + 彻底清除双引号与 Markdown 标记)
# ==============================================================================

COLOR_NAVY_DARK   = RGBColor(27, 58, 92)     # #1B3A5C 主基调·深海蓝
COLOR_TECH_BLUE   = RGBColor(46, 134, 193)   # #2E86C1 辅助·科技蓝
COLOR_AMBER_GOLD  = RGBColor(200, 140, 27)   # #C88C1B 高亮·暖琥珀金
COLOR_EMERALD     = RGBColor(26, 140, 93)    # #1A8C5D 正向·翡翠绿
COLOR_BRICK_RED   = RGBColor(200, 75, 27)    # #C84B1B 警示·砖红
COLOR_CARD_BG     = RGBColor(242, 242, 242)  # #F2F2F2 容器浅灰底
COLOR_CARD_BORDER = RGBColor(218, 224, 230)  # #DAE0E6 卡片微边框
COLOR_CASE_BG     = RGBColor(253, 246, 227)  # #FDF6E3 案例暖杏底
COLOR_TEXT_MAIN   = RGBColor(44, 44, 44)     # #2C2C2C 正文暗炭灰
COLOR_TEXT_MUTED  = RGBColor(107, 107, 107)  # #6B6B6B 次要中灰
COLOR_WHITE       = RGBColor(255, 255, 255)  # #FFFFFF 纯白

FONT_CN = "Microsoft YaHei"   # 严格对标微软雅黑
FONT_EN = "Times New Roman"   # 英文与数字统一 Times New Roman

def sanitize_math_text(text):
    """
    全自动、彻底消除 LaTeX 语法代码与乱码痕迹，转换为符合人类自然阅读直觉的出版级数学符号：
    - 彻底剥离所有 LaTeX 行内/行间数学定界符 $ 与 $$
    - 剥离所有双引号 (", “, ”, \")，提升视觉清爽度
    - 核心符号：\\times -> ×, \\cdot -> ·, \\sum -> ∑, \\Delta -> Δ, \\partial -> ∂, \\int -> ∫
    - 希腊字母：\\alpha, \\beta, \\gamma, \\mu, \\sigma, \\pi, \\lambda, \\rho, \\theta 等
    - 幂次/上标：^2 -> ², ^3 -> ³, ^n -> ⁿ, ^t -> ᵗ, ^-n -> ⁻ⁿ, ^* -> *
    - 下标：_{max} -> max, _{min} -> min, _{10%,3} -> (10%, 3期), _d -> d, _s -> s, _t -> t
    - 分式与根号：\\frac{A}{B} -> (A / B), \\sqrt{A} -> √(A)
    - 箭头/不等号：\\Rightarrow -> ➔, \\rightarrow -> ➔, \\le -> ≤, \\ge -> ≥, \\ne -> ≠, \\approx -> ≈
    """
    if not text:
        return ""
    t = str(text)
    
    # 0. 替换曲线标记中的双撇与双引号
    t = t.replace('D"', "D''").replace('S"', "S''").replace(r'\"', '')
    
    # 1. 剥离所有 LaTeX 数学模式 $ 符号与双引号
    t = t.replace("$", "")
    t = t.replace(r"\text", "")
    t = t.replace('"', '').replace('“', '').replace('”', '')
    
    # 2. 核心数学运算符替换
    t = t.replace(r"\times", "×")
    t = t.replace(r"\cdot", "·")
    t = t.replace(r"\sum", "∑")
    t = t.replace(r"\Delta", "Δ")
    t = t.replace(r"\int", "∫")
    t = t.replace(r"\partial", "∂")
    t = t.replace(r"\alpha", "α")
    t = t.replace(r"\beta", "β")
    t = t.replace(r"\gamma", "γ")
    t = t.replace(r"\epsilon", "ε")
    t = t.replace(r"\mu", "μ")
    t = t.replace(r"\sigma", "σ")
    t = t.replace(r"\pi", "π")
    t = t.replace(r"\Pi", "Π")
    t = t.replace(r"\lambda", "λ")
    t = t.replace(r"\rho", "ρ")
    t = t.replace(r"\theta", "θ")
    t = t.replace(r"\approx", "≈")
    t = t.replace(r"\le", "≤")
    t = t.replace(r"\ge", "≥")
    t = t.replace(r"\ne", "≠")
    t = t.replace(r"\infty", "∞")
    t = t.replace(r"\Rightarrow", "➔")
    t = t.replace(r"\rightarrow", "➔")
    t = t.replace(r"\dots", "…")
    t = t.replace(r"\cdots", "…")
    t = t.replace(r"\quad", " ")
    t = t.replace(r"\qquad", "  ")
    t = t.replace(r"\%", "%")
    t = t.replace(r"\ln", "ln")
    t = t.replace(r"\log", "log")
    
    # 3. 分数替换 \frac{A}{B} -> (A / B)
    t = re.sub(r"\\frac\{([^{}]+)\}\{([^{}]+)\}", r"(\1 / \2)", t)
    t = re.sub(r"\\frac\{([^{}]+)\}\{([^{}]+)\}", r"(\1 / \2)", t)

    # 4. 根号替换 \sqrt{A} -> √(A)
    t = re.sub(r"\\sqrt\{([^{}]+)\}", r"√(\1)", t)

    # 5. 上标/幂次替换
    t = t.replace("^2", "²").replace("^3", "³").replace("^n", "ⁿ").replace("^-n", "⁻ⁿ").replace("^t", "ᵗ").replace("^*", "*")
    t = t.replace("R^2", "R²").replace("R̄²", "R̄²")
    t = t.replace("Q^2", "Q²").replace("Q^3", "Q³")
    t = t.replace("P^b", "Pᵇ").replace("I^c", "Iᶜ").replace("P_o^d", "Pₒᵈ")
    t = re.sub(r"\^\{([0-9a-zA-Z+-]+)\}", r"^\1", t)
    t = re.sub(r"\^([0-9a-zA-Z*]+)", r"\1", t)

    # 6. 下标清理
    t = t.replace("_{max}", " max").replace("_{min}", " min")
    t = re.sub(r"_\{([^{}]+)\}", r"(\1)", t)
    t = t.replace("Q_d", "Qd").replace("Q_s", "Qs").replace("P_b", "Pb").replace("P_s", "Ps")
    t = t.replace("P_m", "Pm").replace("Q_m", "Qm").replace("P_t", "Pt").replace("Q_t", "Qt")
    t = t.replace("P_w", "Pw").replace("Q_w", "Qw").replace("P_f", "Pf").replace("P_c", "Pc")
    t = t.replace("P_r", "Pr").replace("T_e", "Te").replace("E_p", "Ep").replace("N_f", "Nf")

    # 7. 清理残留的 LaTeX 转义反斜杠和花括号
    t = re.sub(r"\\([a-zA-Z]+)", r"\1", t)
    t = t.replace(r"\_", "_")
    t = t.replace("{", "").replace("}", "")

    return t

def add_rich_paragraph(text_frame, raw_text, font_size=12.5, default_color=COLOR_TEXT_MAIN, default_bold=False, space_after=4, font_name=FONT_CN):
    """
    智能富文本段落渲染器：
    1. 全自动消除 LaTeX / $ / 乱码与多余的双引号 (", “, ”, \")
    2. 智能解析 Markdown **粗体** 标记，生成原生 PowerPoint 加粗 Run，彻底消除原生 ** 符号残留
    """
    p = text_frame.add_paragraph()
    p.space_after = Pt(space_after)
    
    # 1. 数学符号清洗与双引号清除
    cleaned = sanitize_math_text(raw_text)
    
    # 2. 解析 **加粗** 标记
    tokens = re.split(r'(\*\*.*?\*\*)', cleaned)
    for token in tokens:
        if not token:
            continue
        run = p.add_run()
        if token.startswith('**') and token.endswith('**') and len(token) >= 4:
            run.text = token[2:-2]
            run.font.bold = True
            run.font.color.rgb = COLOR_NAVY_DARK if default_color == COLOR_TEXT_MAIN else default_color
        else:
            run.text = token
            run.font.bold = default_bold
            run.font.color.rgb = default_color
        run.font.name = font_name
        run.font.size = Pt(font_size)
    return p

def create_deck(book_title, ch_num_en, ch_title, ch_desc, slides_data):
    """
    通用金牌标准课件构建器 (支持矢量图表自动嵌入 + 讲练融合 + 零双引号/零 Markdown 痕迹)
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.500)
    blank_layout = prs.slide_layouts[6]

    def add_header_footer(slide, nav_text, title_text, page_num):
        nav_box = slide.shapes.add_textbox(Inches(0.50), Inches(0.22), Inches(6.00), Inches(0.35))
        tf = nav_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = sanitize_math_text(nav_text)
        p.font.name = FONT_CN
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_TECH_BLUE

        title_box = slide.shapes.add_textbox(Inches(0.50), Inches(0.55), Inches(11.50), Inches(0.60))
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.text = sanitize_math_text(title_text)
        p2.font.name = FONT_CN
        p2.font.size = Pt(26)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_NAVY_DARK

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.50), Inches(1.20), Inches(1.20), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_AMBER_GOLD
        line.line.fill.background()

        footer_box = slide.shapes.add_textbox(Inches(0.50), Inches(7.00), Inches(6.00), Inches(0.35))
        tff = footer_box.text_frame
        tff.margin_left = tff.margin_top = tff.margin_right = tff.margin_bottom = 0
        pf = tff.paragraphs[0]
        pf.text = f"{sanitize_math_text(book_title)} · {sanitize_math_text(ch_title)}"
        pf.font.name = FONT_CN
        pf.font.size = Pt(10)
        pf.font.color.rgb = COLOR_TEXT_MUTED

        num_box = slide.shapes.add_textbox(Inches(12.20), Inches(7.00), Inches(0.80), Inches(0.35))
        tfn = num_box.text_frame
        tfn.margin_left = tfn.margin_top = tfn.margin_right = tfn.margin_bottom = 0
        pn = tfn.paragraphs[0]
        pn.text = f"{page_num:02d}"
        pn.alignment = PP_ALIGN.RIGHT
        pn.font.name = FONT_EN
        pn.font.size = Pt(11)
        pn.font.color.rgb = COLOR_TEXT_MUTED

    def add_card(slide, left, top, width, height, header_text, header_bg=COLOR_NAVY_DARK, card_bg=COLOR_CARD_BG):
        bg_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        bg_card.fill.solid()
        bg_card.fill.fore_color.rgb = card_bg
        bg_card.line.color.rgb = COLOR_CARD_BORDER
        bg_card.line.width = Pt(1)

        h_height = Inches(0.48)
        h_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, h_height)
        h_card.fill.solid()
        h_card.fill.fore_color.rgb = header_bg
        h_card.line.fill.background()

        tf = h_card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.20)
        p = tf.paragraphs[0]
        p.text = sanitize_math_text(header_text)
        p.font.name = FONT_CN
        p.font.size = Pt(15.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        content_box = slide.shapes.add_textbox(
            left + Inches(0.20), 
            top + h_height + Inches(0.12), 
            width - Inches(0.40), 
            height - h_height - Inches(0.22)
        )
        ctf = content_box.text_frame
        ctf.word_wrap = True
        ctf.margin_left = ctf.margin_top = ctf.margin_right = ctf.margin_bottom = 0
        return ctf

    def add_image_card(slide, left, top, width, height, header_text, image_path, header_bg=COLOR_NAVY_DARK, card_bg=COLOR_CARD_BG):
        bg_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        bg_card.fill.solid()
        bg_card.fill.fore_color.rgb = card_bg
        bg_card.line.color.rgb = COLOR_CARD_BORDER
        bg_card.line.width = Pt(1)

        h_height = Inches(0.48)
        h_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, h_height)
        h_card.fill.solid()
        h_card.fill.fore_color.rgb = header_bg
        h_card.line.fill.background()

        tf = h_card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.20)
        p = tf.paragraphs[0]
        p.text = sanitize_math_text(header_text)
        p.font.name = FONT_CN
        p.font.size = Pt(15.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        if os.path.exists(image_path):
            img_left = left + Inches(0.15)
            img_top = top + h_height + Inches(0.10)
            img_w = width - Inches(0.30)
            img_h = height - h_height - Inches(0.20)
            slide.shapes.add_picture(image_path, img_left, img_top, img_w, img_h)

    def add_bottom_takeaway(slide, text, color=COLOR_NAVY_DARK, bg_color=COLOR_CASE_BG, border_color=COLOR_AMBER_GOLD):
        bot_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.50), Inches(6.45), Inches(12.33), Inches(0.45))
        bot_bar.fill.solid()
        bot_bar.fill.fore_color.rgb = bg_color
        bot_bar.line.color.rgb = border_color
        bot_bar.line.width = Pt(1)
        tf = bot_bar.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = sanitize_math_text(text)
        p.font.name = FONT_CN
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color

    # --- 渲染各页面 ---
    for idx, s_info in enumerate(slides_data):
        page_num = idx + 1
        stype = s_info.get("type", "content")

        if stype == "cover":
            s = prs.slides.add_slide(blank_layout)
            top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = COLOR_NAVY_DARK
            top_bar.line.fill.background()

            c_box = s.shapes.add_textbox(Inches(0.80), Inches(1.20), Inches(11.733), Inches(3.20))
            ctf = c_box.text_frame
            ctf.word_wrap = True

            p0 = ctf.paragraphs[0]
            p0.text = ch_num_en
            p0.font.name = FONT_EN
            p0.font.size = Pt(15)
            p0.font.bold = True
            p0.font.color.rgb = COLOR_TECH_BLUE
            p0.space_after = Pt(6)

            p1 = ctf.add_paragraph()
            p1.text = sanitize_math_text(ch_title)
            p1.font.name = FONT_CN
            p1.font.size = Pt(36)
            p1.font.bold = True
            p1.font.color.rgb = COLOR_NAVY_DARK
            p1.space_after = Pt(10)

            p2 = ctf.add_paragraph()
            p2.text = sanitize_math_text(ch_desc)
            p2.font.name = FONT_CN
            p2.font.size = Pt(14)
            p2.font.color.rgb = COLOR_TEXT_MUTED

            cols = s_info.get("cols", [])
            for i, (ctitle, cdesc) in enumerate(cols):
                col_left = Inches(0.80) + i * Inches(2.95)
                card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_left, Inches(4.50), Inches(2.80), Inches(2.10))
                card.fill.solid()
                card.fill.fore_color.rgb = COLOR_CARD_BG
                card.line.color.rgb = COLOR_CARD_BORDER
                tf = card.text_frame
                tf.margin_left = tf.margin_top = Inches(0.20)
                
                p = tf.paragraphs[0]
                p.text = sanitize_math_text(ctitle)
                p.font.name = FONT_CN
                p.font.size = Pt(15)
                p.font.bold = True
                p.font.color.rgb = COLOR_NAVY_DARK
                p.space_after = Pt(6)

                p_desc = tf.add_paragraph()
                p_desc.text = sanitize_math_text(cdesc)
                p_desc.font.name = FONT_CN
                p_desc.font.size = Pt(12.5)
                p_desc.font.color.rgb = COLOR_TEXT_MUTED

        elif stype == "overview":
            s = prs.slides.add_slide(blank_layout)
            add_header_footer(s, s_info["nav"], s_info["title"], page_num)
            cols = s_info.get("cols", [])
            num_cols = len(cols)
            col_w = (12.33 - (num_cols - 1) * 0.15) / num_cols
            for i, (ctitle, citems) in enumerate(cols):
                c_left = Inches(0.50) + i * Inches(col_w + 0.15)
                ctf = add_card(s, c_left, Inches(1.45), Inches(col_w), Inches(4.80), ctitle, 
                               header_bg=COLOR_NAVY_DARK if i%2==0 else COLOR_TECH_BLUE)
                for item in citems:
                    add_rich_paragraph(ctf, item, font_size=12, default_color=COLOR_TEXT_MAIN, space_after=6)
            if "takeaway" in s_info:
                add_bottom_takeaway(s, s_info["takeaway"])

        elif stype == "bridge":
            s = prs.slides.add_slide(blank_layout)
            bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.500))
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
            bg.line.fill.background()

            top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = COLOR_NAVY_DARK
            top_bar.line.fill.background()

            box = s.shapes.add_textbox(Inches(1.20), Inches(2.10), Inches(11.00), Inches(3.50))
            tf = box.text_frame
            tf.word_wrap = True

            p0 = tf.paragraphs[0]
            p0.text = s_info["sec_num"]
            p0.font.name = FONT_EN
            p0.font.size = Pt(20)
            p0.font.bold = True
            p0.font.color.rgb = COLOR_TECH_BLUE
            p0.space_after = Pt(12)

            p1 = tf.add_paragraph()
            p1.text = sanitize_math_text(s_info["sec_title"])
            p1.font.name = FONT_CN
            p1.font.size = Pt(36)
            p1.font.bold = True
            p1.font.color.rgb = COLOR_NAVY_DARK
            p1.space_after = Pt(16)

            line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.20), Inches(3.85), Inches(1.80), Inches(0.05))
            line.fill.solid()
            line.fill.fore_color.rgb = COLOR_AMBER_GOLD
            line.line.fill.background()

            box2 = s.shapes.add_textbox(Inches(1.20), Inches(4.15), Inches(10.50), Inches(1.50))
            tf2 = box2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = sanitize_math_text(s_info["sec_desc"])
            p2.font.name = FONT_CN
            p2.font.size = Pt(15)
            p2.font.color.rgb = COLOR_TEXT_MUTED

            num_box = s.shapes.add_textbox(Inches(12.20), Inches(7.00), Inches(0.80), Inches(0.35))
            tfn = num_box.text_frame
            pn = tfn.paragraphs[0]
            pn.text = f"{page_num:02d}"
            pn.alignment = PP_ALIGN.RIGHT
            pn.font.name = FONT_EN
            pn.font.size = Pt(11)
            pn.font.color.rgb = COLOR_TEXT_MUTED

        elif stype == "2col":
            s = prs.slides.add_slide(blank_layout)
            add_header_footer(s, s_info["nav"], s_info["title"], page_num)
            
            c_bg_l = COLOR_CASE_BG if s_info.get("is_case") else COLOR_CARD_BG
            c_bg_r = COLOR_CASE_BG if s_info.get("is_case") else COLOR_CARD_BG

            # Left side
            left_img = s_info.get("left_image")
            if left_img and os.path.exists(left_img):
                add_image_card(s, Inches(0.50), Inches(1.45), Inches(5.95), Inches(4.80), 
                               s_info["left_title"], left_img, header_bg=COLOR_NAVY_DARK, card_bg=c_bg_l)
            else:
                ctf_l = add_card(s, Inches(0.50), Inches(1.45), Inches(5.95), Inches(4.80), 
                                 s_info["left_title"], header_bg=COLOR_NAVY_DARK, card_bg=c_bg_l)
                for p_item in s_info.get("left_content", []):
                    add_rich_paragraph(ctf_l, p_item, font_size=12.5, default_color=COLOR_TEXT_MAIN, space_after=4)

            # Right Card
            ctf_r = add_card(s, Inches(6.85), Inches(1.45), Inches(5.95), Inches(4.80), 
                             s_info["right_title"], header_bg=COLOR_TECH_BLUE, card_bg=c_bg_r)
            for p_item in s_info.get("right_content", []):
                add_rich_paragraph(ctf_r, p_item, font_size=12.5, default_color=COLOR_TEXT_MAIN, space_after=4)

            if "takeaway" in s_info:
                add_bottom_takeaway(s, s_info["takeaway"])

        elif stype == "summary_cards":
            s = prs.slides.add_slide(blank_layout)
            add_header_footer(s, s_info["nav"], s_info["title"], page_num)
            items = s_info.get("items", [])
            for i, (stitle, sdesc) in enumerate(items):
                top_pos = Inches(1.45) + i * Inches(1.20)
                tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.50), top_pos, Inches(0.90), Inches(1.05))
                tag.fill.solid()
                tag.fill.fore_color.rgb = COLOR_NAVY_DARK
                tag.line.fill.background()
                tf_t = tag.text_frame
                tf_t.vertical_anchor = MSO_ANCHOR.MIDDLE
                p_t = tf_t.paragraphs[0]
                p_t.text = f"{i+1:02d}"
                p_t.font.name = FONT_EN
                p_t.font.size = Pt(20)
                p_t.font.bold = True
                p_t.font.color.rgb = COLOR_WHITE
                p_t.alignment = PP_ALIGN.CENTER

                tbox = s.shapes.add_textbox(Inches(1.55), top_pos, Inches(3.60), Inches(1.05))
                tf_tb = tbox.text_frame
                tf_tb.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf_tb.margin_left = tf_tb.margin_top = tf_tb.margin_right = tf_tb.margin_bottom = 0
                p_tb = tf_tb.paragraphs[0]
                p_tb.text = sanitize_math_text(stitle)
                p_tb.font.name = FONT_CN
                p_tb.font.size = Pt(14)
                p_tb.font.bold = True
                p_tb.font.color.rgb = COLOR_NAVY_DARK

                dbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.30), top_pos, Inches(7.53), Inches(1.05))
                dbox.fill.solid()
                dbox.fill.fore_color.rgb = COLOR_CARD_BG
                dbox.line.color.rgb = COLOR_CARD_BORDER
                tf_d = dbox.text_frame
                tf_d.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf_d.margin_left = tf_d.margin_right = Inches(0.20)
                tf_d.margin_top = tf_d.margin_bottom = Inches(0.10)
                add_rich_paragraph(tf_d, sdesc, font_size=12.5, default_color=COLOR_TEXT_MAIN, space_after=0)

            if "takeaway" in s_info:
                add_bottom_takeaway(s, s_info["takeaway"])

    return prs
