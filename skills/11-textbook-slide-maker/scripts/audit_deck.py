import os
import sys
import re
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

def audit_deck(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"Error: File not found at {pptx_path}")
        return False

    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)
    print(f"=== Auditing Presentation: {os.path.basename(pptx_path)} ===")
    print(f"Total Slides: {total_slides}")

    img_count = 0
    violations = []

    for idx, slide in enumerate(prs.slides):
        slide_imgs = [s for s in slide.shapes if s.shape_type == 13]
        img_count += len(slide_imgs)
        slide_text = ""
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    slide_text += " " + p.text
        for char in ["**", '"', "“", "”", "$", r"\times", r"\frac"]:
            if char in slide_text:
                violations.append((idx + 1, char))

    print(f"Total Embedded Images: {img_count}")
    if violations:
        print("Found violations (Page, Issue):", violations)
        return False
    else:
        print("ALL_PASS: 质检 100% 完美通过！零双引号、零 Markdown 星号、零 LaTeX 乱码！")
        return True

if __name__ == "__main__":
    if len(sys.argv) > 1:
        audit_deck(sys.argv[1])
    else:
        print("Usage: py -3.12 audit_deck.py <path_to_pptx>")
