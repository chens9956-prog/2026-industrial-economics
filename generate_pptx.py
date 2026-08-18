import os
import glob
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation(yaml_path):
    with open(yaml_path, 'r', encoding='utf-8') as f:
        # Load yaml. Note: our YAML starts with comments.
        data = yaml.safe_load(f)
        
    slides_data = data.get('Slides', [])
    
    prs = Presentation()
    # 16:9 ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    MCKINSEY_BLUE = RGBColor(10, 42, 94) # #0A2A5E
    GRAY = RGBColor(80, 80, 80)
    
    for slide_info in slides_data:
        title_text = slide_info.get('Title', '')
        layout_type = slide_info.get('Layout', '')
        content = slide_info.get('Content', {})
        
        # We use a blank layout for maximum control
        slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(slide_layout)
        
        # Add Title at the top
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(40)
        p.font.name = 'Microsoft YaHei'
        p.font.color.rgb = MCKINSEY_BLUE
        
        # Draw a blue accent line under the title
        line = slide.shapes.add_shape(
            9, # msoShapeLine
            Inches(1), Inches(1.5), Inches(11.33), Inches(0)
        )
        line.line.color.rgb = MCKINSEY_BLUE
        line.line.width = Pt(2)
        
        # Render based on Layout type
        if layout_type == 'Hero':
            body_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(2))
            tf_body = body_box.text_frame
            p_body = tf_body.paragraphs[0]
            # Content is a list for Hero
            if isinstance(content, list) and len(content) > 0:
                p_body.text = content[0]
            elif isinstance(content, str):
                p_body.text = content
            p_body.alignment = PP_ALIGN.CENTER
            p_body.font.size = Pt(36)
            p_body.font.bold = True
            p_body.font.color.rgb = GRAY
            p_body.font.name = 'Microsoft YaHei'
            
        elif layout_type == 'Side-by-Side':
            left_text = content.get('Left', '')
            right_text = content.get('Right', '')
            
            # Left box
            box1 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(5), Inches(3))
            p1 = box1.text_frame.paragraphs[0]
            p1.text = left_text
            p1.font.size = Pt(28)
            p1.font.color.rgb = GRAY
            p1.font.name = 'Microsoft YaHei'
            
            # Right box
            box2 = slide.shapes.add_textbox(Inches(7), Inches(2.5), Inches(5), Inches(3))
            p2 = box2.text_frame.paragraphs[0]
            p2.text = right_text
            p2.font.size = Pt(28)
            p2.font.color.rgb = GRAY
            p2.font.name = 'Microsoft YaHei'
            
        elif layout_type == 'Three-column':
            col1 = content.get('Col1', '')
            col2 = content.get('Col2', '')
            col3 = content.get('Col3', '')
            
            w = 3.3
            for i, c_text in enumerate([col1, col2, col3]):
                box = slide.shapes.add_textbox(Inches(1 + i*4), Inches(2.5), Inches(w), Inches(3))
                p = box.text_frame.paragraphs[0]
                p.text = str(i+1) + ". " + c_text
                p.font.size = Pt(24)
                p.font.color.rgb = GRAY
                p.font.name = 'Microsoft YaHei'
                p.font.bold = True
                
        elif layout_type == 'Summary-check':
            box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(4))
            tf = box.text_frame
            tf.word_wrap = True
            if isinstance(content, list):
                for item in content:
                    p = tf.add_paragraph()
                    p.text = "✔️ " + item
                    p.font.size = Pt(24)
                    p.font.color.rgb = GRAY
                    p.font.name = 'Microsoft YaHei'
                    # Spacing
                    p.space_after = Pt(14)
                    
    # Generate output path
    base_name = os.path.basename(yaml_path).replace('.txt', '.pptx').replace('.yaml', '.pptx')
    out_path = os.path.join(os.path.dirname(yaml_path), base_name)
    prs.save(out_path)
    print(f"Generated PPTX: {out_path}")

def main():
    directory = r"I:\4产业经济学\简报施工图"
    # We copied yaml to txt, so they might be .txt or .yaml. We will process .txt files that look like blueprints
    for f in glob.glob(os.path.join(directory, "*.txt")):
        if "施工图" in f:
            try:
                create_presentation(f)
            except Exception as e:
                print(f"Error processing {f}: {e}")
                
if __name__ == "__main__":
    main()
