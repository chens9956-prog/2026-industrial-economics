import os
import sys
import re
import pymupdf
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

if sys.stdout is not None:
    try:
        sys.stdout.reconfigure(encoding='utf-8')
    except Exception:
        pass

class WatermarkRemover:
    def __init__(self, mode="smart", custom_box=None, bg_color_mode="auto", custom_bg_rgb=None):
        """
        :param mode: 'smart' (智能文本/图标定位), 'manual' (按精确像素/比例尺寸定位)
        :param custom_box: dict(width_pt=140, height_pt=26, margin_right_pt=10, margin_bottom_pt=10)
        :param bg_color_mode: 'auto' (智能吸取周边像素), 'white', 'black', 'custom'
        :param custom_bg_rgb: (r, g, b) normalized 0.0-1.0
        """
        self.mode = mode
        self.custom_box = custom_box or {
            "width_pt": 140,
            "height_pt": 26,
            "margin_right_pt": 10,
            "margin_bottom_pt": 10
        }
        self.bg_color_mode = bg_color_mode
        self.custom_bg_rgb = custom_bg_rgb or (1.0, 1.0, 1.0)

    def sample_background_color(self, page, target_rect):
        """
        在目标水印框周边采样真实的背景底色，确保擦除色与背景 100% 融为一体
        """
        try:
            w = page.rect.width
            h = page.rect.height

            # 取样点 1：水印框上方 4pt 处
            sample_x = max(0, min(w - 2, target_rect.x0 + 10))
            sample_y = max(0, min(h - 2, target_rect.y0 - 4))
            
            # 若上方超出，取左侧 4pt 处
            if target_rect.y0 - 4 <= 0:
                sample_x = max(0, min(w - 2, target_rect.x0 - 4))
                sample_y = max(0, min(h - 2, target_rect.y0 + 5))

            pix = page.get_pixmap(clip=pymupdf.Rect(sample_x, sample_y, sample_x + 2, sample_y + 2))
            rgb = pix.pixel(0, 0)
            return (rgb[0] / 255.0, rgb[1] / 255.0, rgb[2] / 255.0)
        except Exception:
            return (1.0, 1.0, 1.0)

    def clean_pdf(self, input_path, output_path=None):
        if output_path is None:
            base, ext = os.path.splitext(input_path)
            output_path = f"{base}_clean{ext}"

        doc = pymupdf.open(input_path)
        total_pages = len(doc)
        cleaned_count = 0

        keywords = [
            "Gemini Notebook", "Gemini", "NotebookLM", "Notebook LM", 
            "Notebook", "Google Notebook"
        ]

        for page_idx in range(total_pages):
            page = doc[page_idx]
            pw = page.rect.width
            ph = page.rect.height

            target_rects = []

            # 1. 尝试基于文本检索精确定位水印 (智能模式)
            if self.mode == "smart":
                for kw in keywords:
                    found_rects = page.search_for(kw)
                    for r in found_rects:
                        # 水印只可能位于右下角区域 (右侧 35%, 底部 25%)
                        if r.x0 >= pw * 0.60 and r.y0 >= ph * 0.75:
                            # 精准微调扩展：向左扩展 26pt (覆盖小图标)，上下各扩 3pt
                            expanded_rect = pymupdf.Rect(
                                max(0, r.x0 - 26),
                                max(0, r.y0 - 3),
                                min(pw, r.x1 + 4),
                                min(ph, r.y1 + 4)
                            )
                            target_rects.append(expanded_rect)
                            break
                    if target_rects:
                        break

            # 2. 如果智能文字搜索未命中，则使用高精度紧凑手动/默认选框
            if not target_rects:
                bw = self.custom_box.get("width_pt", 140)
                bh = self.custom_box.get("height_pt", 26)
                mr = self.custom_box.get("margin_right_pt", 10)
                mb = self.custom_box.get("margin_bottom_pt", 10)

                target_rect = pymupdf.Rect(
                    pw - mr - bw,
                    ph - mb - bh,
                    pw - mr,
                    ph - mb
                )
                target_rects.append(target_rect)

            # 3. 对每个命中区域进行精细色彩匹配擦除
            for r in target_rects:
                if self.bg_color_mode == "auto":
                    fill_color = self.sample_background_color(page, r)
                elif self.bg_color_mode == "white":
                    fill_color = (1.0, 1.0, 1.0)
                elif self.bg_color_mode == "black":
                    fill_color = (0.0, 0.0, 0.0)
                else:
                    fill_color = self.custom_bg_rgb

                # 应用极小 Redact 注释抹除
                page.add_redact_annot(r, fill=fill_color)
                page.apply_redactions()
                cleaned_count += 1

        doc.save(output_path, garbage=4, deflate=True)
        return {
            "success": True,
            "output_path": output_path,
            "total_pages": total_pages,
            "cleaned_count": cleaned_count,
            "file_type": "PDF"
        }

    def clean_pptx(self, input_path, output_path=None):
        """
        PPTX 中直接删除形状对象，实现 0 痕迹残留
        """
        if output_path is None:
            base, ext = os.path.splitext(input_path)
            output_path = f"{base}_clean{ext}"

        prs = Presentation(input_path)
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        removed_count = 0

        def is_watermark(shape):
            try:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # 检查文本
                if shape.has_text_frame:
                    t = shape.text.lower()
                    if any(kw in t for kw in ['gemini notebook', 'notebooklm', 'notebook lm', 'gemini']):
                        return True

                # 检查形状名称
                s_name = shape.name.lower()
                if any(kw in s_name for kw in ['gemini', 'notebook', 'logo', 'watermark']):
                    return True

                # 位置判定：右下角极小形状
                in_bottom_right = (left >= slide_w * 0.70) and (top >= slide_h * 0.78)
                if in_bottom_right and (width < slide_w * 0.25 and height < slide_h * 0.18):
                    return True
            except Exception:
                pass
            return False

        def clean_shapes_in_container(container):
            nonlocal removed_count
            sp_to_delete = [s for s in container.shapes if is_watermark(s)]
            for s in sp_to_delete:
                sp = s._element
                sp.getparent().remove(sp)
                removed_count += 1

        for slide in prs.slides:
            clean_shapes_in_container(slide)

        for master in prs.slide_masters:
            clean_shapes_in_container(master)
            for layout in master.slide_layouts:
                clean_shapes_in_container(layout)

        prs.save(output_path)
        return {
            "success": True,
            "output_path": output_path,
            "removed_count": removed_count,
            "file_type": "PPTX"
        }

    def process_file(self, file_path, output_path=None):
        ext = os.path.splitext(file_path)[1].lower()
        if ext in ['.pptx', '.ppt']:
            return self.clean_pptx(file_path, output_path)
        elif ext == '.pdf':
            return self.clean_pdf(file_path, output_path)
        else:
            raise ValueError(f"不支持的文件格式: {ext}")
